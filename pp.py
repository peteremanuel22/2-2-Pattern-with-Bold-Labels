# pp.py
import io
import os
import gc
import tempfile
import streamlit as st
from PIL import Image, ImageOps
from pptx import Presentation
from pptx.util import Cm, Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from reportlab.lib.pagesizes import A4, landscape
from reportlab.pdfgen import canvas
from reportlab.lib.utils import ImageReader
from reportlab.lib.colors import black

# ---------------------------
# Configuration
# ---------------------------
MAX_FILES = 150

# A4 landscape (cm)
A4_WIDTH_CM = 29.7
A4_HEIGHT_CM = 21.0

# For PDF (points)
A4_LANDSCAPE = landscape(A4)  # ~ (842, 595) points

# Spacing between pictures (inches on ruler)
SPACING_INCH = 0.25  # fixed horizontal and vertical gap between the 4 pictures

# Label band height (fraction of each cell height)
TEXT_FRACTION = 0.12  # 12% of cell height reserved for the label

# Default font settings for the centered label above each picture
LABEL_FONT_SIZE_PT = 24
LABEL_FONT_BOLD = True

# JPEG compression
DEFAULT_JPEG_QUALITY = 50  # you can lower to 35–45 for even smaller memory/size

# Target rendering DPI for slide cells -> pixel sizing during downscale
DEFAULT_TARGET_DPI = 150  # 96 is OK; 150–200 looks crisper with modest memory

# ---------------------------
# Utilities
# ---------------------------
def set_white_background_pptx(slide):
    """Solid white background (safety)."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)

def get_blank_layout(prs: Presentation):
    """Prefer 'Blank' slide layout."""
    for layout in prs.slide_layouts:
        if layout.name.strip().lower() == "blank":
            return layout
    return prs.slide_layouts[0]

def pil_open_transposed(file_obj) -> Image.Image:
    """Open image from UploadedFile/file-like and fix EXIF orientation."""
    img = Image.open(file_obj)
    img = ImageOps.exif_transpose(img)
    return img

def rotate_if_portrait(img: Image.Image) -> Image.Image:
    """Rotate 90° clockwise if portrait."""
    if img.height > img.width:
        return img.rotate(270, expand=True)
    return img

def sanitize_title_from_name(name: str) -> str:
    base = os.path.basename(name)
    stem = os.path.splitext(base)[0]
    return stem.strip() or "Slide"

def inches_from_emu(emu_val: int) -> float:
    """EMU to inches. 1 inch = 914400 EMU."""
    return float(emu_val) / 914400.0

# ---------------------------
# Pattern Geometry (2×2 with fixed 0.25" gaps on a blank slide)
# ---------------------------
def compute_2x2_pattern_geometry_pptx(prs, spacing_inch=SPACING_INCH, text_fraction=TEXT_FRACTION, img_ar=1.5):
    """
    Compute a full-slide, centered 2×2 pattern geometry (no title), preserving AR,
    using fixed gaps. Returns (pattern_left, pattern_top, cell_w, cell_h, text_h, img_area_h, gh, gv).
    All values are EMUs.
    """
    slide_w = prs.slide_width
    slide_h = prs.slide_height

    gh = Inches(spacing_inch)  # horizontal gap in EMUs
    gv = Inches(spacing_inch)  # vertical gap in EMUs

    # cell_width = r_img * (1 - text_fraction) * cell_height
    r_cell = img_ar * (1.0 - text_fraction)
    if r_cell <= 0:
        r_cell = 1e-6

    # Maximize cell height ch subject to slide constraints:
    # 2 * r_cell * ch + gh <= slide_w
    # 2 * ch + gv <= slide_h
    ch_from_width = (slide_w - gh) / (2.0 * r_cell)
    ch_from_height = (slide_h - gv) / 2.0
    ch = min(ch_from_width, ch_from_height)
    if ch <= 0:
        gh = gv = 0
        ch_from_width = slide_w / (2.0 * r_cell)
        ch_from_height = slide_h / 2.0
        ch = max(1, min(ch_from_width, ch_from_height))

    cw = r_cell * ch
    text_h = ch * text_fraction
    img_area_h = ch - text_h

    pattern_w = 2 * cw + gh
    pattern_h = 2 * ch + gv
    pattern_left = (slide_w - pattern_w) / 2.0
    pattern_top = (slide_h - pattern_h) / 2.0

    return (pattern_left, pattern_top, cw, ch, text_h, img_area_h, gh, gv)

# ---------------------------
# Memory-safe image processing
# ---------------------------
def downscale_and_save_temp_jpeg(
    img: Image.Image,
    target_w_inch: float,
    target_h_inch: float,
    dpi: int,
    quality: int,
) -> str:
    """
    Downscale PIL image to (<= target_w_px, <= target_h_px) keeping AR,
    re-encode as JPEG to a NamedTemporaryFile on disk, and return the path.
    """
    target_w_px = max(1, int(round(target_w_inch * dpi)))
    target_h_px = max(1, int(round(target_h_inch * dpi)))

    # Work on a copy to avoid mutating original
    work = img.copy()
    # Downscale in place (keeps AR). Use high-quality resample.
    work.thumbnail((target_w_px, target_h_px), Image.Resampling.LANCZOS)

    # Ensure RGB (remove alpha if any) and write to temp file
    if work.mode in ("RGBA", "LA"):
        bg = Image.new("RGB", work.size, (255, 255, 255))
        alpha = work.split()[-1]
        bg.paste(work.convert("RGB"), mask=alpha)
        work = bg
    else:
        work = work.convert("RGB")

    ntf = tempfile.NamedTemporaryFile(delete=False, suffix=".jpg")
    ntf_path = ntf.name
    ntf.close()
    work.save(ntf_path, format="JPEG", quality=quality, optimize=True)
    work.close()
    return ntf_path

# ---------------------------
# Add 2×2 pattern (PPTX)
# ---------------------------
def add_2x2_pattern_pptx_from_file(
    slide,
    prs,
    temp_jpeg_path: str,
    img_ar: float,
    label_text: str,
    spacing_inch=SPACING_INCH,
    text_fraction=TEXT_FRACTION,
    label_font_pt=LABEL_FONT_SIZE_PT,
    label_bold=LABEL_FONT_BOLD,
):
    pl, pt, cw, ch, text_h, img_area_h, gh, gv = compute_2x2_pattern_geometry_pptx(
        prs, spacing_inch=spacing_inch, text_fraction=text_fraction, img_ar=img_ar
    )

    def cell_origin(r: int, c: int):
        left = pl + (cw + gh) * c
        top = pt + (ch + gv) * r
        return left, top

    # Add 4 cells
    for r in range(2):
        for c in range(2):
            left, top = cell_origin(r, c)

            # Text band (centered)
            tbox = slide.shapes.add_textbox(int(left), int(top), int(cw), int(text_h))
            tf = tbox.text_frame
            tf.clear()
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = label_text
            run.font.size = Pt(label_font_pt)
            run.font.bold = label_bold

            # Image below the label
            img_left = left
            img_top = top + text_h
            pic = slide.shapes.add_picture(temp_jpeg_path, int(img_left), int(img_top), width=int(cw))
            # Center vertically in image area (guard tiny rounding)
            pic.left = int(img_left + (cw - pic.width) / 2.0)
            pic.top = int(img_top + (img_area_h - pic.height) / 2.0)

# ---------------------------
# Add 2×2 pattern (PDF)
# ---------------------------
def add_2x2_pattern_pdf_from_file(
    cpdf: canvas.Canvas,
    page_w: float,
    page_h: float,
    temp_jpeg_path: str,
    img_ar: float,
    label_text: str,
    spacing_inch=SPACING_INCH,
    text_fraction=TEXT_FRACTION,
    label_font_pt=LABEL_FONT_SIZE_PT,
):
    gap_h = 72.0 * spacing_inch
    gap_v = 72.0 * spacing_inch

    r_cell = img_ar * (1.0 - text_fraction)
    if r_cell <= 0:
        r_cell = 1e-6

    ch_from_width = (page_w - gap_h) / (2.0 * r_cell)
    ch_from_height = (page_h - gap_v) / 2.0
    ch = min(ch_from_width, ch_from_height)
    if ch <= 0:
        gap_h = gap_v = 0
        ch_from_width = page_w / (2.0 * r_cell)
        ch_from_height = page_h / 2.0
        ch = max(1, min(ch_from_width, ch_from_height))

    cw = r_cell * ch
    text_h = ch * text_fraction
    img_area_h = ch - text_h

    pattern_w = 2 * cw + gap_h
    pattern_h = 2 * ch + gap_v
    pattern_left = (page_w - pattern_w) / 2.0
    pattern_bottom = (page_h - pattern_h) / 2.0

    img_reader = ImageReader(temp_jpeg_path)

    cpdf.setFont("Helvetica-Bold", label_font_pt)
    cpdf.setFillColor(black)

    for r in range(2):
        for col in range(2):
            x = pattern_left + (cw + gap_h) * col
            y = pattern_bottom + (ch + gap_v) * r

            # Text center
            text_cx = x + cw / 2.0
            text_cy = y + ch - (text_h / 2.0)
            cpdf.drawCentredString(text_cx, text_cy, label_text)

            # Image area bottom-left
            cpdf.drawImage(
                img_reader,
                x,
                y,
                width=cw,
                height=img_area_h,
                preserveAspectRatio=True,
                anchor='sw',
            )

# ---------------------------
# Streamlit App
# ---------------------------
st.set_page_config(page_title="A4 Landscape — 2×2 Pattern (PPTX/PDF) — Memory‑Safe", page_icon="📄", layout="centered")
st.title("📄 A4 Landscape — 2×2 Pattern with Bold Labels (PPTX or PDF)")
st.caption(
    "Efficient, memory‑safe processing: images are handled one‑by‑one, downscaled to cell size, "
    "re‑encoded to JPEG, then immediately released. This avoids Cloud resource limit crashes."
)

uploaded_files = st.file_uploader(
    "Upload 1 to 150 images",
    type=["jpg", "jpeg", "png", "bmp", "tif", "tiff", "webp"],
    accept_multiple_files=True,
    help="Images are processed sequentially to keep memory low."
)

col0, col1, col2, col3 = st.columns([1,1,1,1])
with col0:
    output_format = st.radio("Output", ["PPTX", "PDF"], horizontal=True)
with col1:
    jpeg_quality = st.slider("JPEG quality", 30, 90, value=DEFAULT_JPEG_QUALITY, step=5)
with col2:
    target_dpi = st.slider("Target DPI", 96, 220, value=DEFAULT_TARGET_DPI, step=6, help="Used to downscale into the cell size")
with col3:
    gap_inch = st.select_slider("Gap (inches)", options=[0.10, 0.15, 0.20, 0.25, 0.30, 0.40], value=SPACING_INCH)

default_name_pptx = "images_to_a4_landscape_2x2_pattern_memsafe.pptx"
default_name_pdf  = "images_to_a4_landscape_2x2_pattern_memsafe.pdf"
out_name = st.text_input("Output filename", value=(default_name_pptx if output_format == "PPTX" else default_name_pdf))

generate = st.button("Generate")

if generate:
    if not uploaded_files:
        st.error("Please upload at least 1 image.")
    elif len(uploaded_files) > MAX_FILES:
        st.error(f"Please upload at most {MAX_FILES} images.")
    else:
        try:
            progress = st.progress(0)
            status = st.empty()

            if output_format == "PPTX":
                prs = Presentation()
                prs.slide_width = Cm(A4_WIDTH_CM)
                prs.slide_height = Cm(A4_HEIGHT_CM)
                blank_layout = get_blank_layout(prs)

                # Precompute pattern geometry once per *typical* AR; we still compute per-image below for exact AR
                for idx, up in enumerate(uploaded_files, start=1):
                    up.seek(0)
                    # Open & normalize on the fly
                    with pil_open_transposed(up) as img0:
                        img = rotate_if_portrait(img0)
                        img_ar = img.width / img.height if img.height else 1.0

                        # Compute final cell sizes to know target pixel bounds
                        pl, pt, cw, ch, text_h, img_area_h, gh, gv = compute_2x2_pattern_geometry_pptx(
                            prs, spacing_inch=gap_inch, text_fraction=TEXT_FRACTION, img_ar=img_ar
                        )
                        cell_w_in = inches_from_emu(cw)
                        img_h_in = inches_from_emu(img_area_h)

                        # Downscale + encode to temp JPEG on disk (reused 4x on this slide)
                        temp_jpeg = downscale_and_save_temp_jpeg(
                            img, target_w_inch=cell_w_in, target_h_inch=img_h_in, dpi=target_dpi, quality=jpeg_quality
                        )

                    label_text = sanitize_title_from_name(up.name)

                    # Build slide
                    slide = prs.slides.add_slide(blank_layout)
                    set_white_background_pptx(slide)

                    add_2x2_pattern_pptx_from_file(
                        slide=slide,
                        prs=prs,
                        temp_jpeg_path=temp_jpeg,
                        img_ar=img_ar,
                        label_text=label_text,
                        spacing_inch=gap_inch,
                        text_fraction=TEXT_FRACTION,
                        label_font_pt=LABEL_FONT_SIZE_PT,
                        label_bold=LABEL_FONT_BOLD,
                    )

                    # Cleanup temp + Python refs
                    try:
                        os.unlink(temp_jpeg)
                    except Exception:
                        pass
                    gc.collect()

                    if idx % 3 == 0 or idx == len(uploaded_files):
                        progress.progress(idx / len(uploaded_files))
                        status.text(f"[PPTX] Processed {idx}/{len(uploaded_files)}: {up.name}")

                out_buf = io.BytesIO()
                prs.save(out_buf)
                out_buf.seek(0)
                st.success("PPTX generated successfully!")
                st.download_button(
                    label="⬇️ Download PPTX",
                    data=out_buf,
                    file_name=(out_name.strip() or default_name_pptx),
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                )

            else:
                page_w, page_h = A4_LANDSCAPE
                out_buf = io.BytesIO()
                cpdf = canvas.Canvas(out_buf, pagesize=A4_LANDSCAPE)

                for idx, up in enumerate(uploaded_files, start=1):
                    up.seek(0)
                    with pil_open_transposed(up) as img0:
                        img = rotate_if_portrait(img0)
                        img_ar = img.width / img.height if img.height else 1.0

                        # Compute cell sizes to determine downscale target in inches (PDF uses points; convert)
                        # For consistency with PPTX geometry, compute with same method in "inches"
                        # Derive cw,ch via same formula:
                        # r_cell = img_ar*(1-text_fraction)
                        r_cell = img_ar * (1.0 - TEXT_FRACTION) if img_ar > 0 else 1.0
                        gap_h = 72.0 * gap_inch
                        gap_v = 72.0 * gap_inch
                        ch_from_width = (page_w - gap_h) / (2.0 * r_cell)
                        ch_from_height = (page_h - gap_v) / 2.0
                        ch = max(1.0, min(ch_from_width, ch_from_height))
                        cw = r_cell * ch
                        text_h = ch * TEXT_FRACTION
                        img_area_h = ch - text_h

                        # Convert cw/img_area_h points -> inches (1 inch = 72 pt)
                        cell_w_in = cw / 72.0
                        img_h_in = img_area_h / 72.0

                        temp_jpeg = downscale_and_save_temp_jpeg(
                            img, target_w_inch=cell_w_in, target_h_inch=img_h_in, dpi=target_dpi, quality=jpeg_quality
                        )

                    label_text = sanitize_title_from_name(up.name)

                    # White page background
                    cpdf.setFillColorRGB(1, 1, 1)
                    cpdf.rect(0, 0, page_w, page_h, fill=1, stroke=0)

                    add_2x2_pattern_pdf_from_file(
                        cpdf, page_w, page_h, temp_jpeg, img_ar, label_text,
                        spacing_inch=gap_inch, text_fraction=TEXT_FRACTION, label_font_pt=LABEL_FONT_SIZE_PT
                    )

                    cpdf.showPage()

                    try:
                        os.unlink(temp_jpeg)
                    except Exception:
                        pass
                    gc.collect()

                    if idx % 3 == 0 or idx == len(uploaded_files):
                        progress.progress(idx / len(uploaded_files))
                        status.text(f"[PDF] Processed {idx}/{len(uploaded_files)}: {up.name}")

                cpdf.save()
                out_buf.seek(0)
                st.success("PDF generated successfully!")
                st.download_button(
                    label="⬇️ Download PDF",
                    data=out_buf,
                    file_name=(out_name.strip() or default_name_pdf),
                    mime="application/pdf",
                )

        except Exception as e:
            st.error(f"An error occurred: {e}")
